from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


BG = RGBColor(248, 241, 230)
CARD = RGBColor(255, 250, 242)
INK = RGBColor(30, 26, 23)
MUTED = RGBColor(107, 91, 80)
GOLD = RGBColor(245, 201, 76)
LINE = RGBColor(232, 221, 206)


SLIDES = [
    (
        "MaRe",
        [
            "One luxury scalp-health app with guest, internal, salon-owner, and client experiences",
            "Hackathon solution across Flutter, Flask, backend, AI agent, and AWS storage layers",
        ],
    ),
    (
        "The Problem",
        [
            "MaRe wins when the experience feels exclusive, diagnostic, and personal.",
            "National growth usually pushes brands toward generic outreach and low-quality content.",
            "That creates the personalized scale paradox: scale makes luxury feel cheap.",
        ],
    ),
    (
        "The Opportunity",
        [
            "Scalp wellness and head spa routines are growing search and service categories.",
            "Luxury salons need ancillary revenue and stronger retail conversion.",
            "MaRe can become the system behind scalp diagnostics, education, and retail lift.",
        ],
    ),
    (
        "Our Solution",
        [
            "Guest mode for public discovery.",
            "One sign-in for all MaRe roles.",
            "Role picker for internal, salon-owner, and client dashboards.",
            "Shared AI, storage, and review systems under one app shell.",
        ],
    ),
    (
        "Flutter App",
        [
            "One MaRe shell across iOS, Android, and Web.",
            "Shows guest mode, sign-in, role selection, and role-specific dashboards.",
            "Includes json_annotation models, a local backend stub, and agent fallback rules.",
        ],
    ),
    (
        "Web App",
        [
            "Simple Flask + Python app for the lightweight web version.",
            "Includes Flask routes, a Vercel-compatible Python entrypoint, and AWS S3-oriented image storage.",
            "Matches the Flutter business story so judges can review it instantly in browser.",
        ],
    ),
    (
        "Why It Wins",
        [
            "Protects brand fidelity while increasing prospecting speed.",
            "Improves outreach quality without sounding like AI spam.",
            "Scales content volume for AI search and social channels.",
            "Turns AI failure into a safe, visible review workflow.",
        ],
    ),
    (
        "Architecture",
        [
            "Frontend + backend + AI agent in both solution paths.",
            "Shared domain model: prospect, outreach, content, review, AI error, and image storage.",
            "Yellow dot means AI-controlled or fallback-enabled surface.",
        ],
    ),
    (
        "ROI Story",
        [
            "Better prospect targeting for luxury-fit salons.",
            "Higher retail conversion through the MaRe Eye consultation flow.",
            "More qualified partner conversations and more content throughput.",
            "Stronger AI-search visibility on scalp-health and wellness queries.",
        ],
    ),
    (
        "Roadmap",
        [
            "Connect live salon datasets and verification providers.",
            "Integrate real multimodal generation and fully signed AWS uploads for images and video assets.",
            "Add CRM history, analytics, permissions, and production observability.",
            "Deploy regional expansion playbooks nationwide.",
        ],
    ),
]


def add_text_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    accent = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10.9), Inches(0.35), Inches(0.42), Inches(0.42)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.6), Inches(11.6), Inches(5.9)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(9.5), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    paragraph = title_frame.paragraphs[0]
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = INK

    body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(10.4), Inches(3.9))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    for index, bullet in enumerate(bullets):
        paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = MUTED
        paragraph.space_after = Pt(12)
        paragraph.bullet = True


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for title, bullets in SLIDES:
        add_text_slide(prs, title, bullets)

    output = Path(__file__).with_name("MaRe_Luxury_Growth_Engine.pptx")
    prs.save(output)
    print(output)


if __name__ == "__main__":
    main()
